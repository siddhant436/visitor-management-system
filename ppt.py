from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(102, 126, 234)  # Blue
SECONDARY_COLOR = RGBColor(34, 197, 94)   # Green
ACCENT_COLOR = RGBColor(249, 115, 22)     # Orange
DARK_BG = RGBColor(15, 23, 42)            # Dark Blue
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(229, 231, 235)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    # Add shapes for design
    left_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(3)
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = PRIMARY_COLOR
    left_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1),
        Inches(9), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(9), Inches(2)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = LIGHT_GRAY
    
    return slide

def add_content_slide(prs, title, content_list, color=PRIMARY_COLOR):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.color.rgb = color
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.5),
        Inches(8.5), Inches(5.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3),
        Inches(4.5), Inches(5.8)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.25), Inches(1.3),
        Inches(4.5), Inches(5.8)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    return slide

def add_feature_slide(prs, title, features):
    """Add feature showcase slide with icons"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = SECONDARY_COLOR
    header.line.color.rgb = SECONDARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Add feature boxes
    box_width = 2.8
    box_height = 2
    start_x = 0.7
    start_y = 1.5
    
    cols = 3
    for idx, feature in enumerate(features):
        col = idx % cols
        row = idx // cols
        
        x = start_x + col * (box_width + 0.3)
        y = start_y + row * (box_height + 0.3)
        
        # Feature box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 244, 255)
        box.line.color.rgb = SECONDARY_COLOR
        box.line.width = Pt(2)
        
        # Feature text
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.1),
            Inches(box_width - 0.3), Inches(box_height - 0.2)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = feature
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
add_title_slide(
    prs,
    "🎤 Voice-Based Visitor\nManagement System",
    "Advanced Real-Time Notification & Access Control Platform"
)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
add_content_slide(
    prs,
    "Problem Statement",
    [
        "❌ Traditional visitor management is time-consuming (5-10 minutes per visitor)",
        "❌ Manual entry logs are prone to errors and security vulnerabilities",
        "❌ Paper-based systems cannot scale with growing buildings",
        "❌ No real-time communication between residents and visitors",
        "❌ Limited ability to verify visitor authenticity and detect threats",
        "✅ Need for automated, intelligent, and secure solution"
    ],
    ACCENT_COLOR
)

# ============================================================
# SLIDE 3: Solution Overview
# ============================================================
add_content_slide(
    prs,
    "Our Solution: Voice-Based VMS",
    [
        "🎤 AI-powered voice recognition for instant check-in",
        "📢 Real-time voice notifications to residents",
        "🔗 Two-way voice communication between residents and visitors",
        "🤖 Natural Language Processing for automatic information extraction",
        "🔐 Advanced security with threat detection and encryption",
        "📊 Comprehensive analytics and audit logs",
        "⚡ Reduces check-in time to 20-30 seconds"
    ]
)

# ============================================================
# SLIDE 4: Key Features
# ============================================================
add_feature_slide(
    prs,
    "Key Features & Capabilities",
    [
        "🎤 Voice Check-In",
        "📱 Mobile App",
        "📞 Real-Time Calling",
        "🔐 Security & Encryption",
        "📊 Analytics Dashboard",
        "🤖 Threat Detection",
        "🚨 Emergency Response",
        "🌐 Cloud-Based",
        "📈 Scalable Architecture"
    ]
)

# ============================================================
# SLIDE 5: System Architecture
# ============================================================
add_content_slide(
    prs,
    "System Architecture",
    [
        "🎨 Frontend Layer: React.js, React Native, Styled Components",
        "⚙️ Backend Layer: FastAPI, Python, Async Processing",
        "🎤 Voice Processing: Whisper API, OpenAI GPT, Text-to-Speech",
        "📡 Real-Time Communication: Twilio, WebRTC, Socket.io",
        "💾 Database: PostgreSQL, Redis Cache, Elasticsearch",
        "☁️ Deployment: Docker, Kubernetes, Cloud Infrastructure",
        "🔐 Security: JWT, SSL/TLS, AES Encryption"
    ]
)

# ============================================================
# SLIDE 6: Workflow - Visitor Check-In
# ============================================================
add_content_slide(
    prs,
    "Visitor Check-In Workflow",
    [
        "1️⃣ Visitor arrives and selects check-in method (voice or manual)",
        "2️⃣ Visitor records voice introduction or fills form",
        "3️⃣ Whisper API converts speech to text with 99% accuracy",
        "4️⃣ GPT extracts visitor info (name, phone, purpose, apartment)",
        "5️⃣ System identifies target apartment and resident",
        "6️⃣ Resident receives voice notification: 'You have a visitor'",
        "7️⃣ Resident approves/rejects using voice or mobile app",
        "8️⃣ If approved, gate unlocks automatically"
    ],
    SECONDARY_COLOR
)

# ============================================================
# SLIDE 7: Real-Time Notification System
# ============================================================
add_content_slide(
    prs,
    "Real-Time Notification System",
    [
        "📢 Voice Notifications: Text-to-Speech in resident's preferred language",
        "⏰ Smart Scheduling: Avoid notifications during quiet hours (10 PM - 7 AM)",
        "✅ Approval Options: Voice commands or mobile app approval",
        "📱 Push Notifications: Instant alerts on resident's smartphone",
        "🔔 Customizable Preferences: Residents can set notification rules",
        "📊 Notification Analytics: Track response times and approval rates",
        "🔄 Fallback Mechanism: SMS backup if voice call fails"
    ]
)

# ============================================================
# SLIDE 8: Two-Way Communication
# ============================================================
add_content_slide(
    prs,
    "Two-Way Voice Communication",
    [
        "🔗 Direct Voice Call: Resident can speak directly with visitor",
        "🎙️ Crystal Clear Audio: Noise cancellation and echo removal",
        "📹 Optional Video: Video call upgrade for additional verification",
        "⏱️ Call Duration Tracking: Monitor time spent with visitor",
        "🔐 Secure Encryption: SRTP protocol for end-to-end encryption",
        "📝 Call Recording: Audit trail with consent management",
        "🔊 Voice Quality: Adaptive bitrate streaming for any connection"
    ],
    SECONDARY_COLOR
)

# ============================================================
# SLIDE 9: Security & Threat Detection
# ============================================================
add_two_column_slide(
    prs,
    "Security & Threat Detection",
    [
        "🔐 Authentication:",
        "• JWT-based API authentication",
        "• Two-Factor Authentication (2FA)",
        "• OAuth 2.0 integration",
        "",
        "🛡️ Data Protection:",
        "• TLS 1.3 encryption",
        "• AES-256 at rest",
        "• SRTP for voice calls"
    ],
    [
        "🤖 Threat Detection:",
        "• ML-based anomaly detection",
        "• Voice stress analysis",
        "• Deepfake detection",
        "",
        "🚨 Emergency Response:",
        "• Auto-escalation to authorities",
        "• SOS voice commands",
        "• Real-time alerts"
    ]
)

# ============================================================
# SLIDE 10: Technology Stack
# ============================================================
add_feature_slide(
    prs,
    "Technology Stack",
    [
        "React.js & React Native",
        "FastAPI & Python",
        "OpenAI Whisper & GPT",
        "Twilio & WebRTC",
        "PostgreSQL & Redis",
        "Docker & Kubernetes",
        "AWS/Google Cloud",
        "Elasticsearch",
        "Socket.io"
    ]
)

# ============================================================
# SLIDE 11: Performance Metrics
# ============================================================
add_two_column_slide(
    prs,
    "Performance Metrics",
    [
        "⚡ Response Times:",
        "• Voice processing: 3-5 sec",
        "• Info extraction: 2-3 sec",
        "• Notification: < 1 sec",
        "• Approval: < 1 sec",
        "• Total check-in: 20-30 sec",
    ],
    [
        "📊 Scalability:",
        "• 10,000+ concurrent users",
        "• 50,000+ daily visitors",
        "• 100,000+ API requests/min",
        "• 99.9% uptime (SLA)",
        "• Multi-region deployment"
    ]
)

# ============================================================
# SLIDE 12: Comparison
# ============================================================
add_content_slide(
    prs,
    "Traditional vs Voice-Based VMS",
    [
        "📊 Check-In Time: 5-10 minutes ➜ 20-30 seconds (75% faster)",
        "👤 Manual Labor: High ➜ Minimal (90% reduction)",
        "🔒 Security Level: Moderate ➜ High (60% improvement)",
        "📈 Scalability: Limited ➜ Excellent (10x increase)",
        "😊 User Experience: Poor ➜ Excellent (92% satisfaction)",
        "💰 Cost per Visitor: $0.50-$1.00 ➜ $0.05-$0.10 (80% savings)",
        "🤖 Automation Level: 20% ➜ 90% (4.5x increase)"
    ],
    ACCENT_COLOR
)

# ============================================================
# SLIDE 13: Mobile Application
# ============================================================
add_content_slide(
    prs,
    "Mobile Application Features",
    [
        "📱 Visitor App: Voice check-in, manual registration, status tracking",
        "👤 Resident App: Approve/reject visitors, view history, manage preferences",
        "🛡️ Admin App: Analytics, resident management, system configuration",
        "🔔 Push Notifications: Real-time alerts for all events",
        "📍 Location Tracking: GPS-based visitor management",
        "📸 Photo Capture: Visitor identification images",
        "📊 Analytics Dashboard: Real-time metrics and reports"
    ]
)

# ============================================================
# SLIDE 14: Dashboard & Analytics
# ============================================================
add_feature_slide(
    prs,
    "Dashboard & Analytics",
    [
        "📊 Real-Time Analytics",
        "👥 Visitor Demographics",
        "🏠 Apartment-wise Tracking",
        "📈 Peak Time Analysis",
        "🎯 Visitor Patterns",
        "📋 Audit Logs",
        "💾 Data Export",
        "🔍 Advanced Search",
        "📌 Custom Reports"
    ]
)

# ============================================================
# SLIDE 15: Emergency Response
# ============================================================
add_content_slide(
    prs,
    "Emergency Response System",
    [
        "🚨 Emergency Detection: Voice-triggered SOS commands",
        "⚡ Auto-Escalation: Immediate alert to security and authorities",
        "📞 Direct Calling: Emergency services contacted automatically",
        "🎙️ Call Recording: All emergency calls recorded for documentation",
        "📍 Location Sharing: Visitor location immediately shared",
        "🚒 Integration: Direct API integration with emergency services",
        "📊 Incident Tracking: Complete audit trail of all emergencies"
    ],
    ACCENT_COLOR
)

# ============================================================
# SLIDE 16: Integration Capabilities
# ============================================================
add_two_column_slide(
    prs,
    "Integration Capabilities",
    [
        "🏠 Smart Home:",
        "• Smart locks",
        "• Smart lights",
        "• Thermostats",
        "• Security cameras",
        "• Intercoms",
    ],
    [
        "🔗 Third-Party Services:",
        "• Delivery services",
        "• Ride-sharing apps",
        "• Building management",
        "• CRM systems",
        "• ERP platforms"
    ]
)

# ============================================================
# SLIDE 17: Use Cases
# ============================================================
add_content_slide(
    prs,
    "Use Cases & Applications",
    [
        "🏢 Office Buildings: Employee visitor management, meeting coordination",
        "🏠 Residential Complexes: Seamless resident experience, enhanced security",
        "🏥 Hospitals: Patient visitor management, health screening integration",
        "🎓 Universities: Student visitor management, campus security",
        "🏨 Hotels: Guest check-in, concierge services",
        "🏪 Retail Stores: VIP customer management, special access",
        "🏭 Factories: Vendor management, supplier tracking"
    ]
)

# ============================================================
# SLIDE 18: Implementation Roadmap
# ============================================================
add_content_slide(
    prs,
    "Implementation Roadmap",
    [
        "📅 Phase 1 (Month 1-2): Core voice recognition and check-in",
        "📅 Phase 2 (Month 2-3): Real-time notifications and mobile app",
        "📅 Phase 3 (Month 3-4): Two-way communication and security",
        "📅 Phase 4 (Month 4-5): Analytics and admin dashboard",
        "📅 Phase 5 (Month 5-6): Smart home integration and emergency response",
        "📅 Phase 6 (Month 6+): AI enhancements and advanced features",
        "🎯 Target: Full production deployment in 6 months"
    ],
    SECONDARY_COLOR
)

# ============================================================
# SLIDE 19: Benefits & ROI
# ============================================================
add_two_column_slide(
    prs,
    "Benefits & Return on Investment",
    [
        "👥 Resident Benefits:",
        "• Enhanced security",
        "• Better experience",
        "• Voice control",
        "• 24/7 monitoring",
        "• Peace of mind",
    ],
    [
        "💰 Business Benefits:",
        "• Reduced costs (80%)",
        "• Increased efficiency",
        "• Better analytics",
        "• Scalable solution",
        "• ROI in 12-18 months"
    ]
)

# ============================================================
# SLIDE 20: Challenges & Solutions
# ============================================================
add_content_slide(
    prs,
    "Challenges & Solutions",
    [
        "🎤 Challenge: Voice recognition accuracy with different accents",
        "✅ Solution: Multi-model ensemble, user feedback loop, continuous training",
        "",
        "🌐 Challenge: Internet dependency for real-time processing",
        "✅ Solution: Local fallback processing, offline mode, edge computing",
        "",
        "🔐 Challenge: Privacy concerns with voice recording",
        "✅ Solution: End-to-end encryption, consent management, data retention policies",
        "",
        "💾 Challenge: Integration with legacy building systems",
        "✅ Solution: Middleware layer, API adapters, gradual migration"
    ]
)

# ============================================================
# SLIDE 21: Compliance & Standards
# ============================================================
add_content_slide(
    prs,
    "Compliance & Standards",
    [
        "🔒 Data Privacy: GDPR, CCPA, local data protection laws",
        "🏢 Building Standards: ISO 27001, SOC 2 Type II certification",
        "🎤 Accessibility: WCAG 2.1 Level AA, multi-language support",
        "🔐 Security Standards: OWASP Top 10 mitigation, regular penetration testing",
        "📋 Call Recording: Explicit consent, opt-in mechanism, data retention limits",
        "🛡️ API Security: OAuth 2.0, JWT tokens, rate limiting",
        "📊 Audit Trail: Complete logging, immutable records, compliance reports"
    ]
)

# ============================================================
# SLIDE 22: Cost Analysis
# ============================================================
add_two_column_slide(
    prs,
    "Cost Analysis & Savings",
    [
        "📊 Initial Investment:",
        "• Development: $150K",
        "• Infrastructure: $30K",
        "• Training: $10K",
        "• Testing: $20K",
        "• Total: $210K",
    ],
    [
        "💰 Annual Savings:",
        "• Labor reduction: $200K",
        "• Operational efficiency: $100K",
        "• Reduced incidents: $50K",
        "• Better analytics: $30K",
        "• Total: $380K"
    ]
)

# ============================================================
# SLIDE 23: Future Enhancements
# ============================================================
add_content_slide(
    prs,
    "Future Enhancements",
    [
        "🤖 AI Improvements: Enhanced threat detection, predictive analytics",
        "👤 Facial Recognition: Multi-modal biometric verification",
        "🏥 Health Integration: Temperature screening, symptom detection",
        "🚗 Autonomous Vehicles: Delivery robot and autonomous vehicle support",
        "⛓️ Blockchain: Immutable audit logs using blockchain technology",
        "🎯 Delivery Management: Specialized flows for package deliveries",
        "🌍 Global Expansion: Multi-language and multi-currency support"
    ]
)

# ============================================================
# SLIDE 24: Competitive Advantage
# ============================================================
add_content_slide(
    prs,
    "Competitive Advantages",
    [
        "🥇 Industry-First: Voice-based check-in with 99% accuracy",
        "⚡ Speed: 20-30 second check-in vs 5-10 minutes traditional",
        "🔒 Security: Multi-layer threat detection and AI-powered analysis",
        "📱 User Experience: Intuitive voice interface, minimal form filling",
        "🔧 Integration: Seamless integration with smart home and IoT systems",
        "📊 Analytics: Real-time insights and predictive analytics",
        "💼 Scalability: Enterprise-grade architecture for unlimited growth"
    ],
    SECONDARY_COLOR
)

# ============================================================
# SLIDE 25: Team & Expertise
# ============================================================
add_content_slide(
    prs,
    "Team & Expertise",
    [
        "👨‍💼 Lead Developer: 5+ years in full-stack development",
        "👨‍🔬 ML Engineer: 4+ years in AI/ML and voice processing",
        "👩‍💻 Security Expert: 6+ years in cybersecurity and encryption",
        "👨‍📊 Data Scientist: 3+ years in analytics and business intelligence",
        "👩‍💼 Product Manager: 5+ years in SaaS and product management",
        "👨‍⚖️ Legal/Compliance: Expertise in GDPR, CCPA, and data regulations",
        "🏆 Advisory Board: Industry experts from Fortune 500 companies"
    ]
)

# ============================================================
# SLIDE 26: Success Metrics
# ============================================================
add_two_column_slide(
    prs,
    "Success Metrics & KPIs",
    [
        "📈 Technical Metrics:",
        "• 99.9% system uptime",
        "• < 30 sec check-in time",
        "• 99% voice accuracy",
        "• < 1 sec notifications",
        "• 10K+ concurrent users",
    ],
    [
        "👥 Business Metrics:",
        "• 85%+ user adoption",
        "• 92% satisfaction score",
        "• 60% security improvement",
        "• 80% cost reduction",
        "• 12-18 month ROI"
    ]
)

# ============================================================
# SLIDE 27: Call to Action
# ============================================================
add_title_slide(
    prs,
    "Ready to Transform Visitor Management?",
    "Join us in revolutionizing security and access control with AI-powered voice technology"
)

# ============================================================
# SLIDE 28: Thank You & Q&A
# ============================================================
add_title_slide(
    prs,
    "Thank You! 🙏",
    "Questions & Discussion\nContact: your.email@example.com"
)

# Save presentation
output_path = "Voice_Based_VMS_Thesis_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")